from docx import Document
from pptx import Presentation
from tests.ooxml_fixtures import make_docx, make_pptx


def test_make_docx_roundtrips(tmp_path):
    path = make_docx(tmp_path / "s.docx")
    doc = Document(str(path))
    assert any("Claire Martin" in p.text for p in doc.paragraphs)


def test_make_pptx_roundtrips(tmp_path):
    path = make_pptx(tmp_path / "s.pptx")
    prs = Presentation(str(path))
    texts = [sh.text_frame.text for sl in prs.slides for sh in sl.shapes
             if sh.has_text_frame]
    assert any("Claire Martin" in t for t in texts)
