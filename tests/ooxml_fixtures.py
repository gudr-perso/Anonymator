from pathlib import Path
from docx import Document
from pptx import Presentation
from pptx.util import Inches


def make_docx(path: Path) -> Path:
    """Document Word couvrant corps, tableau, en-tête, pied, zone de texte et
    métadonnées — tous conservés par le round-trip python-docx.

    Note : les commentaires / notes de bas de page vivent dans des parties
    séparées (comments.xml, footnotes.xml) qui, dans un vrai .docx, sont
    référencées par des relations et donc conservées à la sauvegarde. On ne
    les injecte pas ici (python-docx les abandonnerait faute de relation) : le
    mécanisme de post-passe est validé isolément dans test_ooxml_xml_parts.py.
    """
    doc = Document()
    doc.core_properties.author = "Alice Durand"
    doc.core_properties.last_modified_by = "Bob Martin"

    # Corps, avec un run scindé manuellement pour tester le remap.
    p = doc.add_paragraph("Contact : ")
    p.add_run("Claire ")
    p.add_run("Mar")
    p.add_run("tin")

    # Tableau
    table = doc.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Fournisseur"
    table.rows[0].cells[1].text = "Claire Martin"

    # En-tête / pied
    section = doc.sections[0]
    section.header.paragraphs[0].text = "Dossier de Claire Martin"
    section.footer.paragraphs[0].text = "Rédigé par Claire Martin"

    doc.save(str(path))
    _inject_textbox(path)
    return path


def _inject_textbox(path: Path) -> None:
    """Ajoute au document.xml une zone de texte (txbxContent) contenant
    « Claire Martin » — que python-docx ne crée pas directement, mais qu'il
    conserve car elle réside dans la partie principale."""
    import zipfile
    from lxml import etree
    W = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    with zipfile.ZipFile(path) as z:
        names = z.namelist()
        data = {n: z.read(n) for n in names}

    doc_root = etree.fromstring(data["word/document.xml"])
    body = doc_root.find(f"{{{W}}}body")
    txbx = etree.SubElement(body, f"{{{W}}}txbxContent")
    txbx.append(etree.fromstring(
        f'<w:p xmlns:w="{W}"><w:r><w:t>Zone de Claire Martin</w:t></w:r></w:p>'))
    data["word/document.xml"] = etree.tostring(doc_root, xml_declaration=True,
                                               encoding="UTF-8", standalone=True)

    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        for n in names:
            z.writestr(n, data[n])


def make_pptx(path: Path) -> Path:
    """Présentation couvrant une zone de texte, un tableau, un groupe et
    les notes du présentateur."""
    prs = Presentation()
    prs.core_properties.author = "Alice Durand"
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Client Claire Martin"

    table = slide.shapes.add_table(1, 2, Inches(1), Inches(3),
                                   Inches(4), Inches(1)).table
    table.cell(0, 0).text = "Nom"
    table.cell(0, 1).text = "Claire Martin"

    notes = slide.notes_slide.notes_text_frame
    notes.text = "Présenté par Claire Martin"

    prs.save(str(path))
    return path
