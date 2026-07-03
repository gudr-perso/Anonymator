from datetime import datetime
from pptx import Presentation
from anonymator.referential import Referential
from anonymator.ner import FakeNer
from anonymator.files.ooxml import pptx_io
from tests.ooxml_fixtures import make_pptx


def _anonymize(tmp_path):
    src = make_pptx(tmp_path / "src.pptx")
    ref = Referential.load_default()
    ner = FakeNer({"Claire Martin": "PERSON"})
    out, report = pptx_io.anonymize_document(
        src, ner, ref, tmp_path, when=datetime(2026, 7, 3, 9, 0, 0))
    return src, out, report


def test_output_name_and_textbox_masked(tmp_path):
    _, out, _ = _anonymize(tmp_path)
    assert out.name == "src_ano_20260703090000.pptx"
    prs = Presentation(str(out))
    texts = [sh.text_frame.text for sl in prs.slides for sh in sl.shapes
             if sh.has_text_frame]
    assert "Client [PERSONNE]" in texts


def test_table_and_notes_masked(tmp_path):
    _, out, _ = _anonymize(tmp_path)
    prs = Presentation(str(out))
    slide = prs.slides[0]
    table = next(sh.table for sh in slide.shapes if sh.has_table)
    assert table.cell(0, 1).text == "[PERSONNE]"
    assert "[PERSONNE]" in slide.notes_slide.notes_text_frame.text


def test_metadata_purged(tmp_path):
    _, out, _ = _anonymize(tmp_path)
    prs = Presentation(str(out))
    assert (prs.core_properties.author or "") == ""


def test_report_present(tmp_path):
    _, _, report = _anonymize(tmp_path)
    assert any(r["original"] == "Claire Martin" for r in report.to_rows())
