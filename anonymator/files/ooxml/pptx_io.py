from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from anonymator.output_naming import anonymized_path
from anonymator.report.audit import AuditReport
from anonymator.files.ooxml import scan, xml_parts
from anonymator.files.ooxml.text_unit import TextUnit


def _iter_frame(text_frame, location: str):
    for p in text_frame.paragraphs:
        if p.runs:
            yield TextUnit(list(p.runs), location)


def _iter_shapes(shapes, prefix: str):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes, prefix)
            continue
        # has_table n'existe que sur GraphicFrame ; has_text_frame varie selon
        # le type de forme → accès défensif via getattr.
        if getattr(shape, "has_text_frame", False):
            yield from _iter_frame(shape.text_frame, prefix)
        if getattr(shape, "has_table", False):
            for ri, row in enumerate(shape.table.rows, 1):
                for ci, cell in enumerate(row.cells, 1):
                    yield from _iter_frame(
                        cell.text_frame, f"{prefix} / Tableau L{ri}C{ci}")


def iter_main_units(prs):
    for si, slide in enumerate(prs.slides, 1):
        yield from _iter_shapes(slide.shapes, f"Slide {si}")
        if slide.has_notes_slide:
            tf = slide.notes_slide.notes_text_frame
            if tf is not None:
                yield from _iter_frame(tf, f"Slide {si} / Notes")


def anonymize_document(path: Path, ner, ref, output_dir: Path,
                       when: datetime) -> tuple[Path, AuditReport]:
    prs = Presentation(str(path))
    units = list(iter_main_units(prs))
    retained = scan.confirmed_only(scan.scan_units(units, ner, ref))
    report = scan.apply_units(units, retained, ref)
    out = anonymized_path(path, output_dir, when)
    prs.save(str(out))
    xml_parts.postprocess_metadata(out, report)
    return out, report
